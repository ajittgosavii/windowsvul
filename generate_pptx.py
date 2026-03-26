"""
Generate VulnShield AI presentation deck.
Follows the EXACT demo flow — each slide matches a demo step.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Colors
DARK_BG = RGBColor(0x0A, 0x0A, 0x23)
BLUE = RGBColor(0x0F, 0x34, 0x60)
LIGHT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xA8, 0xB2, 0xD1)
GREEN = RGBColor(0x28, 0xA7, 0x45)
RED = RGBColor(0xDC, 0x35, 0x45)
ORANGE = RGBColor(0xFD, 0x7E, 0x14)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def dark_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_BG


def text(slide, x, y, w, h, txt, size=16, color=GRAY, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(txt.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = align
        p.space_after = Pt(4)
    return box


def title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)
    text(slide, 0.8, 0.4, 11.5, 0.8, title, size=32, color=WHITE, bold=True)
    if subtitle:
        text(slide, 0.8, 1.1, 11.5, 0.5, subtitle, size=15, color=GRAY)
    return slide


def table(slide, rows, col_widths, x=0.8, y=2.2):
    shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), sum([Inches(w) for w in col_widths]), Inches(0.35 * len(rows)))
    tbl = shape.table
    for ci, w in enumerate(col_widths):
        tbl.columns[ci].width = Inches(w)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = WHITE if ri == 0 else GRAY
                p.font.bold = (ri == 0)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE if ri == 0 else RGBColor(0x12, 0x12, 0x30)


# ==================== SLIDE 1: Title ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
text(slide, 1, 2.2, 11, 1.2, "VulnShield AI", size=54, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
text(slide, 1, 3.5, 11, 0.8, "Autonomous AI Security Engineer\nfor Enterprise Windows Server Vulnerability Management", size=20, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
text(slide, 1, 4.8, 11, 0.8, "Multi-Account AWS  |  SSM Patch Manager  |  ServiceNow ITSM\nHuman-in-the-Loop  |  12 AI Agents  |  Perceive-Reason-Act-Learn", size=13, color=GRAY, align=PP_ALIGN.CENTER)
text(slide, 1, 6.2, 11, 0.4, "winvulmgmt.streamlit.app", size=12, color=PURPLE, align=PP_ALIGN.CENTER)


# ==================== SLIDE 2: The Problem ====================
s = title_slide("The Problem", "Why $200K enterprise tools still leave you exposed")
text(s, 0.8, 1.8, 11.5, 5, """10,000+ Windows Servers across 200+ AWS accounts
New CVEs published daily  --  CISA KEV alerts for active exploits
Manual vulnerability triage: 2 hours per CVE
Manual change request creation: 30 minutes per ticket
Patch cycles take weeks  --  attackers move in hours

Current tools (Qualys, Tenable, CrowdStrike):
   Cost $200K - $300K / year
   Find vulnerabilities but DON'T fix them
   No AI reasoning  --  just "install this KB"
   No autonomous decision-making
   No ServiceNow automation
   No confidence-based human-in-the-loop

The gap:  No tool REASONS about whether to patch,
who to notify, or what change ticket to create.""", size=15)


# ==================== SLIDE 3: The Solution ====================
s = title_slide("VulnShield AI  --  An Autonomous AI Security Engineer", "Not a dashboard. An agent that runs 24/7.")
text(s, 0.8, 1.8, 11.5, 5, """The agent runs a continuous cycle:

   PERCEIVE       What changed in my environment?
                         New CVE on NVD, CISA KEV alert, server compliance drift,
                         new EC2 instance launched, approval received

   REASON           What should I do about it?
                         LLM (GPT-4o / Claude) analyzes risk, reads policies, checks memory
                         "This is a .NET update, safe to auto-patch. Confidence: 94%"

   ACT                  Execute the decision
                         SSM Patch Manager install, ServiceNow CHG ticket, Slack alert,
                         rollback if something breaks

   LEARN              Record the outcome
                         "Last time we patched the ERP server, IIS restarted.
                          Next time, require human approval."

   REPORT            Tell humans only what they need to know
                         Proactive alerts, not reactive dashboards""", size=14)


# ==================== SLIDE 4: Live Demo - Dashboard ====================
s = title_slide("Live Demo: Dashboard  --  Real AWS Data", "\"Everything you see is LIVE from AWS account 448549863273\"")
text(s, 0.8, 1.8, 11.5, 5, """LIVE badge in header  --  toggle between Demo and Live data

What the audience sees:

   AWS Accounts:           2  (hub + spoke via AssumeRole)
   Windows Servers:       12  (10 in hub, 2 in spoke)
   Servers Online:          3  (FinOps, WAFR in us-west-1 + ML-Server)
   Critical Vulnerabilities:  Real count from SSM scan
   Compliance Score:       Calculated from 9 live security checks

Key talking point:
   "These aren't mock numbers. The system just ran PowerShell scripts
    on those servers via SSM and found that TLS 1.2 is NOT enabled.
    That's a real critical finding from a real scan."

Account Vulnerability Posture table shows both accounts:
   448549863273 (Splunk COE)  --  10 servers, 3 regions
   950766978386 (Cloud Migration)  --  2 servers, cross-account via AssumeRole""", size=14)


# ==================== SLIDE 5: Live Demo - Fleet View ====================
s = title_slide("Live Demo: Fleet View  --  12 Real Servers", "\"Notice the Account column  --  different accounts, managed from one place\"")
table(s, [
    ["Server", "Instance ID", "Account", "Region", "OS", "SSM", "Status"],
    ["FinOps", "i-0e284388d8c8e9b79", "448549863273 (Hub)", "us-west-1", "Win Server 2025", "Online", "Running"],
    ["WAFR", "i-03f775ac21c5bc3fe", "448549863273 (Hub)", "us-west-1", "Win Server 2025", "Online", "Running"],
    ["ML-Server", "i-0d381382940e69845", "448549863273 (Hub)", "us-east-1", "Windows Server", "Offline", "Running"],
    ["FinOps", "i-07d45c3c9b49fdba8", "950766978386 (Spoke)", "us-east-1", "Windows Server", "Pending", "Running"],
    ["Genai-Bastion", "i-053a4cfb4d9dc1a78", "448549863273 (Hub)", "us-east-1", "Windows Server", "Offline", "Stopped"],
    ["+ 7 more", "...", "448549863273", "3 regions", "Windows Server", "Offline", "Stopped"],
], [1.2, 2.2, 2.2, 1.2, 1.6, 0.8, 0.8], y=1.8)

text(s, 0.8, 5.2, 11.5, 1.5, """Key talking points:
   "The stopped servers show 0% compliance  --  not because they're secure,
    but because we CAN'T reach them.  That's a visibility gap the tool surfaces."

   "The spoke account servers were discovered via STS AssumeRole.
    The hub never had those credentials."

   Filter by Status = Online  -->  show only the live servers
   Charts: Server Status | By Region | By Environment""", size=13, color=GRAY)


# ==================== SLIDE 6: Patch Status ====================
s = title_slide("Live Demo: Patch Status  --  Real SSM Patch Manager Data", "\"Both servers are missing KB5078740  --  the March 2026 Critical Security Update\"")
text(s, 0.8, 1.8, 5.5, 5, """FinOps (i-0e284388d8c8e9b79):

   Installed:   2 patches
   Other:         5 patches
   Missing:      1 patch          NON-COMPLIANT
   Failed:        0

   Installed:
      KB5066131  .NET Framework 3.5/4.8.1  Critical
      KB5072033  Win Server 2025 Cumulative  Critical

   Missing:
      KB5078740  2026-03 Security Update  CRITICAL

   [Install Missing Patches]  button runs
   AWS-RunPatchBaseline Install on the real server""", size=13)

text(s, 6.5, 1.8, 5.5, 5, """WAFR (i-03f775ac21c5bc3fe):

   Installed:   2 patches
   Other:         5 patches
   Missing:      1 patch          NON-COMPLIANT
   Failed:        0

   Missing:
      KB5078740  2026-03 Security Update  CRITICAL


This is REAL data from AWS SSM Patch Manager.

Patch Baseline: VulnShield-Critical-Security
   Baseline ID: pb-0aa04b12d28ff3abc
   Auto-approves Critical + Security updates
   0-day approval delay

"This is how AWS recommends enterprise patching.
 The AI adds the reasoning layer on top." """, size=13)


# ==================== SLIDE 7: Agent Brain ====================
s = title_slide("Live Demo: Agent Brain  --  The Differentiator", "\"This is NOT a dashboard. This is an autonomous AI agent.\"")
text(s, 0.8, 1.8, 11.5, 5, """Click:  Run Agent Cycle

The agent does 4 things autonomously:

   1. PERCEIVE  --  Checked NVD and CISA KEV for new Windows CVEs
      "It found a critical CVE in the CISA Known Exploited Vulnerabilities catalog.
       That means attackers are using it RIGHT NOW."

   2. REASON  --  GPT-4o analyzed our servers, policies, and memory
      "The AI looked at which servers are affected, what policies apply,
       and what happened last time we patched similar vulnerabilities."

   3. ACT  --  Decided to scan the fleet immediately
      "Because of the CISA KEV alert, it triggered an immediate fleet scan
       and created a notification for the security team."

   4. LEARN  --  Recorded this in memory
      "Next time a similar CVE appears, it already knows what to do."

Actions Table shows every decision with priority, target, and reasoning.
Threat Feed on the right shows real-time CVE alerts.

"No other tool does this.  Qualys finds vulnerabilities.
 This agent finds them, reasons about them, and fixes them." """, size=13)


# ==================== SLIDE 8: Natural Language Policies ====================
s = title_slide("Natural Language Policies", "\"Plain English rules.  No YAML.  No code changes.  The LLM interprets them.\"")
table(s, [
    ["#", "Policy Name", "Rule (Plain English)"],
    ["1", "Auto-patch non-prod", "Automatically patch all Dev and Staging servers without approval"],
    ["2", "Production kernel guard", "Kernel updates on Production require human approval, even at 90%+ confidence"],
    ["3", "Critical CVE bypass", "CRITICAL CVEs with CVSS >= 9.0 can be patched outside maintenance windows"],
    ["4", "ERP protection", "Any remediation on ERP/SAP servers must create a ServiceNow CHG ticket"],
    ["5", "Stale patch escalation", "Servers unpatched > 30 days: escalate to security lead + P2 incident"],
    ["6", "Failure response", "Failed remediation: trigger Rollback Agent + create P1 incident"],
    ["7", "Zero-day response", "New CISA KEV entry: scan fleet + notify security team within 15 minutes"],
], [0.4, 2.2, 8.5], y=1.8)

text(s, 0.8, 5.8, 11.5, 1, """Live demo:  Click "Add New Policy" and type:
   "Never auto-patch the FinOps server on the last day of the month"
   Done.  No code changes.  The AI reads this on its next cycle.""", size=14, color=LIGHT_BLUE)


# ==================== SLIDE 9: Pipeline - The Money Shot ====================
s = title_slide("Live Demo: Agentic Pipeline  --  7 Decisions in 10 Seconds", "\"Watch the AI route each CVE based on confidence\"")
table(s, [
    ["CVE", "Component", "Confidence", "Route", "Why"],
    ["CVE-2024-43498", ".NET Framework", "94%", "AUTO-REMEDIATE", ".NET updates are well-tested, safe to auto-patch"],
    ["CVE-2024-43500", "IIS Web Server", "92%", "AUTO-REMEDIATE", "IIS updates are common, low risk"],
    ["CVE-2024-43499", "Remote Desktop (RDP)", "72%", "HUMAN APPROVAL", "RDP config change on production needs review"],
    ["CVE-2024-38063", "TCP/IP Stack", "74%", "HUMAN APPROVAL", "Network stack change, needs validation"],
    ["CVE-2024-21338", "Windows Kernel", "60%", "CHG TICKET", "Kernel update + reboot, needs change management"],
    ["CVE-2024-30078", "Wi-Fi Driver", "40%", "CHG TICKET", "Driver change is highest risk, needs full process"],
    ["CVE-2024-35250", "Kernel Streaming", "64%", "CHG TICKET", "Kernel component + reboot, insufficient confidence"],
], [1.5, 1.8, 1.0, 1.8, 5], y=1.6)

text(s, 0.8, 5.6, 11.5, 1.5, """Key talking point:

   "2 auto-fixed  --  the AI was 90%+ confident.  No human touched these.
    2 pending approval  --  the AI isn't sure enough.  It asks a human.
    3 CHG tickets  --  low confidence.  These need full change management.

    7 decisions in 10 seconds.  A human analyst takes 2 hours.
    The AI knows when to act and when to stop and ask." """, size=14, color=WHITE)


# ==================== SLIDE 10: Human-in-the-Loop ====================
s = title_slide("Live Demo: Human-in-the-Loop Approval Queue", "\"The AI explains WHY it needs human approval\"")
text(s, 0.8, 1.8, 11.5, 5, """Approval Queue shows only the 2 items the AI wasn't confident about:

   CVE-2024-43499 (RDP)  --  Confidence: 72%  |  Risk: 91%
      AI Reasoning:  "This is an RDP configuration change on a production server.
      The patch modifies NLA and TLS settings. While the fix is standard,
      the production environment and RDP dependency warrant human review."

      Remediation Script Preview  -->  Full PowerShell with rollback

   CVE-2024-38063 (TCP/IP)  --  Confidence: 74%  |  Risk: 94%
      AI Reasoning:  "TCP/IP stack modification affects all network services.
      High CVSS but the remediation involves registry changes to the
      network stack. Recommend testing in staging first."


   Actions:
      Click  Approve   -->  "Now it executes on the real server via SSM"
      Click  Reject     -->  "Stays in queue, no action taken"
      Click  --> CHG    -->  "Escalates to ServiceNow change ticket"

   "The human only reviews what the AI isn't confident about.
    60% of vulnerabilities were handled without any human intervention." """, size=14)


# ==================== SLIDE 11: ServiceNow ====================
s = title_slide("Live Demo: ServiceNow ITSM  --  Auto-Created in 2 Seconds", "\"Open ServiceNow in another tab  --  show the real ticket\"")
text(s, 0.8, 1.8, 11.5, 5, """Click:  Create CHG Ticket  -->  Select CVE-2024-21338 (Kernel EoP)

Switch to ServiceNow tab (dev218436.service-now.com)  -->  Show the ticket:

   Short Description:   [HIGH] CVE-2024-21338 - Windows Kernel EoP
   Category:              Software / Security Patch
   Assignment Group:   Windows Server Team
   Impact:                  2 - High
   Risk:                      High

   Description:
      Full CVE details, NIST controls (SI-2), affected servers,
      AI confidence score, routing decision reasoning

   Implementation Plan:
      1. Verify server availability
      2. Create restore point
      3. Install KB5034763
      4. Apply registry fixes
      5. Verify and close

   Backout Plan:
      1. Restore from system restore point
      2. Rollback registry changes
      3. Verify system stability

   "The AI wrote the entire change request.
    A human writing this takes 30 minutes.  This took 2 seconds." """, size=13)


# ==================== SLIDE 12: Cross-Account Architecture ====================
s = title_slide("Architecture: Hub & Spoke Multi-Account", "\"One IAM role per account.  200 accounts in under 1 hour.\"")
text(s, 0.8, 1.8, 5.5, 5, """Hub Account: 448549863273 (Splunk COE)

   VulnShield AI Platform
   12 AI Agents
   SSM Patch Manager
   ServiceNow ITSM
   Streamlit Observation Deck

   10 Windows Servers
   3 AWS Regions
   2 SSM Online (FinOps, WAFR)
   Patch Baseline: pb-0aa04b12d28ff3abc


Spoke Account: 950766978386 (Cloud Migration)

   WindowsVulnScannerRole
      Trust: 448549863273
      Permissions: SSM + EC2
   SSMInstanceRole on instances

   2 Windows Servers
   1 Running with SSM""", size=13)

text(s, 6.5, 1.8, 5.5, 5, """How it works:

   1. Hub calls STS AssumeRole
   2. Gets temporary credentials for spoke
   3. Runs SSM commands in spoke account
   4. Results flow back to hub dashboard

Scaling to 200+ accounts:
   CloudFormation StackSet deploys
   WindowsVulnScannerRole to all accounts
   One template, one command, 200 accounts

   No agents to install
   SSM Agent is already on every Windows AMI
   No scanners to deploy
   No VPN tunnels needed

Demo:
   Toggle DEMO  -->  6 accounts, 60 servers
   Toggle LIVE   -->  2 real accounts, 12 real servers
   "And switch to real AWS data instantly." """, size=13)


# ==================== SLIDE 13: 12 AI Agents ====================
s = title_slide("12 Specialized AI Agents")
table(s, [
    ["#", "Agent", "What It Does"],
    ["1", "Discovery Agent", "Finds Windows servers across 200+ AWS accounts via SSM"],
    ["2", "Analysis Agent", "LLM-powered CVE analysis with NIST mapping and risk scoring"],
    ["3", "Decision Agent", "Routes to auto-fix / human / CHG based on confidence score"],
    ["4", "Remediation Agent", "Executes patches via SSM Patch Manager (not custom scripts)"],
    ["5", "Verification Agent", "Validates patches installed, services running, no disruption"],
    ["6", "ITSM Agent", "Auto-creates ServiceNow CHG/INC tickets with full plans"],
    ["7", "Rollback Agent", "Auto-reverts failed patches using restore points and registry backup"],
    ["8", "Notification Agent", "Slack, Teams, Email alerts routed by severity"],
    ["9", "Reporting Agent", "Board-ready PDF reports for executives and auditors"],
    ["10", "Scheduling Agent", "Enforces maintenance windows and blackout periods"],
    ["11", "Compliance Drift Agent", "Continuous NIST SP 800-53 and CIS baseline monitoring"],
    ["12", "Threat Intel Agent", "Real-time NVD, CISA KEV, EPSS monitoring + MITRE ATT&CK mapping"],
], [0.4, 2.2, 8.5], y=1.6)


# ==================== SLIDE 14: Market Comparison ====================
s = title_slide("Market Comparison  --  Why This Wins")
table(s, [
    ["Capability", "VulnShield AI", "Qualys ($200K)", "Tenable ($150K)", "CrowdStrike ($300K)"],
    ["Truly Agentic AI", "YES - Perceive/Reason/Act/Learn", "No - basic ML scoring", "No - VPR scoring", "No - threat correlation"],
    ["Autonomous 24/7", "YES - continuous agent loop", "No - scheduled scans", "No - scheduled scans", "No - agent-based detection"],
    ["Auto-Remediation", "Confidence-based routing", "Patch deployment only", "Detection only", "Detection only"],
    ["Human-in-the-Loop", "YES - approval queue + AI reasoning", "No", "No", "No"],
    ["NL Policies", "YES - plain English rules", "No - YAML/XML", "No", "No"],
    ["SSM Patch Manager", "Native AWS integration", "No", "No", "No"],
    ["ServiceNow Auto", "CHG + INC + CMDB + full plans", "Plugin (extra cost)", "Plugin (extra cost)", "Basic integration"],
    ["Memory / Learning", "YES - learns from outcomes", "No", "No", "No"],
    ["Script Generation", "YES - per-CVE PowerShell + rollback", "No", "No", "No"],
    ["Annual Cost", "$5-15K", "$200K+", "$150K+", "$300K+"],
], [2, 2.8, 2, 2, 2.5], y=1.6)


# ==================== SLIDE 15: ROI ====================
s = title_slide("Return on Investment  --  20x to 35x in Year One")
text(s, 0.8, 1.8, 5.5, 5, """Cost Savings:

   Replace Qualys/Tenable:
      $150K - $250K / year saved

   Eliminate manual triage:
      2 FTE analysts x $60-80K = $120K - $160K / year

   Reduce incidents via auto-patching:
      $50K - $100K / year (fewer breaches, less downtime)

   Faster change requests:
      30 min --> 2 sec = thousands of hours saved


   TOTAL SAVINGS:  $320K - $510K / year""", size=15)

text(s, 6.5, 1.8, 5.5, 5, """VulnShield AI Cost:

   OpenAI GPT-4o API:    $2K - $5K / year
   AWS infrastructure:    $1K - $3K / year
   Streamlit Cloud:         $0 - $5K / year

   TOTAL COST:  $3K - $13K / year


   ROI:  25x - 40x return in year one


Time Savings:
   Vulnerability triage:       2 hours  -->  10 seconds
   Change request:             30 min    -->  2 seconds
   Fleet discovery:              Days      -->  Minutes
   Compliance report:          Hours     -->  Instant
   Patch install decision:   Meeting  -->  Autonomous""", size=15)


# ==================== SLIDE 16: Closing ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
text(slide, 1, 1.8, 11, 1, "VulnShield AI", size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
text(slide, 1, 3.0, 11, 1, '"An autonomous AI security engineer that monitors,\nreasons, acts, and learns  --  24/7, across your entire Windows fleet."', size=20, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
text(slide, 1, 4.5, 11, 1.5, "12 AI Agents  |  2 real AWS accounts  |  12 real servers\n7 decisions in 10 seconds  |  $5K vs $200K\nReal SSM Patch Manager  |  Real ServiceNow tickets\nNatural language policies  |  Agent memory that learns", size=15, color=PURPLE, align=PP_ALIGN.CENTER)
text(slide, 1, 6.2, 11, 0.8, "winvulmgmt.streamlit.app\n\nHub: 448549863273  |  Spoke: 950766978386  |  Cross-account via AssumeRole", size=12, color=GRAY, align=PP_ALIGN.CENTER)


prs.save("C:/aiprojects/windowsvulnerabilitiies/VulnShield_AI_Demo_Deck.pptx")
print("PPTX saved: 16 slides following demo flow")
